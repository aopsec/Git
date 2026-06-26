# -*- coding: utf-8 -*-
"""OpenBox v0.2.0 — professional presentation deck (recap of all project snapshots)."""
import pathlib
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

HERE = pathlib.Path(__file__).parent
FIGS = HERE / "figs"

# ---- palette: midnight navy dominant · blue/green support · gold accent ----
NAVY   = RGBColor(0x0C, 0x22, 0x33)
NAVY2  = RGBColor(0x12, 0x31, 0x49)
INK    = RGBColor(0x1A, 0x23, 0x2E)
BLUE   = RGBColor(0x2E, 0x6F, 0xA8)
GREEN  = RGBColor(0x1F, 0x9E, 0x75)
BORDO  = RGBColor(0x8A, 0x1F, 0x3D)
GOLD   = RGBColor(0xE8, 0xA8, 0x19)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ICE    = RGBColor(0xCA, 0xDC, 0xFC)
MUTE   = RGBColor(0x6B, 0x7A, 0x8C)
MUTEDK = RGBColor(0x90, 0xA4, 0xBA)
LINE   = RGBColor(0xD4, 0xDA, 0xE2)
TINT   = RGBColor(0xF1, 0xF5, 0xF9)

SW, SH = 13.333, 7.5
prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

IMG = {p.name: Image.open(p).size for p in FIGS.glob("*.png")}

def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = bg
    return s

def tb(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, sp=2, leading=None):
    """runs: list of paragraphs; each paragraph is list of (text,size,color,bold,italic) run-tuples or a single tuple."""
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = wrap; tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    paras = runs if isinstance(runs[0], list) else [[r] for r in runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp); p.space_before = Pt(0)
        if leading: p.line_spacing = leading
        for (text, size, color, bold, italic) in para:
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
            r.font.color.rgb = color; r.font.name = "Calibri"
    return box

def circle(s, cx, cy, d, fill, line=None, lw=1.0):
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx-d/2), Inches(cy-d/2), Inches(d), Inches(d))
    o.fill.solid(); o.fill.fore_color.rgb = fill
    if line: o.line.color.rgb = line; o.line.width = Pt(lw)
    else: o.line.fill.background()
    o.shadow.inherit = False
    return o

def rrect(s, x, y, w, h, fill, line=None, lw=1.0, radius=0.06):
    r = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: r.fill.background()
    else: r.fill.solid(); r.fill.fore_color.rgb = fill
    if line: r.line.color.rgb = line; r.line.width = Pt(lw)
    else: r.line.fill.background()
    r.shadow.inherit = False
    try: r.adjustments[0] = radius
    except Exception: pass
    return r

def pill(s, x, y, w, h, fill, text, color, size=11.5, bold=True):
    p = rrect(s, x, y, w, h, fill, radius=0.5)
    tf = p.text_frame; tf.word_wrap = False
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    pa = tf.paragraphs[0]; pa.alignment = PP_ALIGN.CENTER
    r = pa.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = "Calibri"
    return p

def header(s, num, category, title, dark=False):
    """section badge motif (gold circle + number) + category + title."""
    tcol = WHITE if dark else INK
    circle(s, 1.02, 0.92, 0.62, GOLD)
    tb(s, 0.72, 0.70, 0.6, 0.45, [[(num, 16, NAVY, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    tb(s, 1.55, 0.60, 10.5, 0.32, [[(category.upper(), 12, GOLD, True, False)]], wrap=False)
    t = tb(s, 1.52, 0.88, 11.0, 0.7, [[(title, 27, tcol, True, False)]], wrap=False)
    t.text_frame.paragraphs[0].runs[0].font.name = "Cambria"
    return 1.95  # content top y

def footer(s, idx, total=15, dark=False):
    col = MUTEDK if dark else MUTE
    tb(s, 0.72, 7.06, 6.0, 0.3, [[("OpenBox v0.2.0  ·  UniCEUB · Introdução à Computação · 2026", 9, col, False, False)]], wrap=False)
    tb(s, 11.3, 7.06, 1.7, 0.3, [[(f"{idx:02d} / {total:02d}", 9, col, False, False)]], align=PP_ALIGN.RIGHT, wrap=False)

def fit(imgname, boxw, boxh):
    iw, ih = Image.open(FIGS/imgname).size; r = min(boxw/iw, boxh/ih)
    return iw*r, ih*r

def framed_pic(s, imgname, x, y, boxw, boxh, dark=False, pad=0.0, center=True):
    """fit image inside box; draw a frame (white card on dark, thin border on light)."""
    fw, fh = fit(imgname, boxw-2*pad, boxh-2*pad)
    ix = x + (boxw-fw)/2 if center else x+pad
    iy = y + (boxh-fh)/2
    if dark:
        rrect(s, ix-0.14, iy-0.14, fw+0.28, fh+0.28, WHITE, radius=0.05)
    else:
        rrect(s, ix-0.10, iy-0.10, fw+0.20, fh+0.20, WHITE, line=LINE, lw=1.0, radius=0.05)
    s.shapes.add_picture(str(FIGS/imgname), Inches(ix), Inches(iy), Inches(fw), Inches(fh))
    return ix, iy, fw, fh

def bullets(s, x, y, w, h, items, color=INK, size=13, accent=GREEN, gap=4, leading=1.05):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(gap); p.space_before = Pt(0); p.line_spacing = leading
        rb = p.add_run(); rb.text = "▸ "; rb.font.size = Pt(size); rb.font.bold = True; rb.font.color.rgb = accent; rb.font.name="Calibri"
        rt = p.add_run(); rt.text = it; rt.font.size = Pt(size); rt.font.color.rgb = color; rt.font.name="Calibri"
    return box

def stat(s, x, y, w, number, label, ncol=GREEN, lcol=MUTE, nsize=34):
    t = tb(s, x, y, w, 0.7, [[(number, nsize, ncol, True, False)]], align=PP_ALIGN.LEFT, wrap=False)
    t.text_frame.paragraphs[0].runs[0].font.name = "Cambria"
    tb(s, x, y+0.62, w, 0.6, [[(label, 11, lcol, False, False)]])

def table(s, x, y, w, headers, rows, colw, hsize=11, bsize=10.5, rowh=0.34, hfill=NAVY, hcol=WHITE):
    nrows = len(rows)+1; ncols = len(headers)
    gt = s.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(rowh*nrows)).table
    gt.first_row = False; gt.horz_banding = False
    for j, cw in enumerate(colw): gt.columns[j].width = Inches(cw)
    for j, htxt in enumerate(headers):
        c = gt.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = hfill
        c.margin_left=Inches(0.08); c.margin_right=Inches(0.05); c.margin_top=Inches(0.03); c.margin_bottom=Inches(0.03)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]; r = p.add_run(); r.text = htxt
        r.font.size = Pt(hsize); r.font.bold = True; r.font.color.rgb = hcol; r.font.name="Calibri"
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = gt.cell(i+1, j); c.fill.solid(); c.fill.fore_color.rgb = WHITE if i%2==0 else TINT
            c.margin_left=Inches(0.08); c.margin_right=Inches(0.05); c.margin_top=Inches(0.02); c.margin_bottom=Inches(0.02)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]; r = p.add_run(); r.text = val
            r.font.size = Pt(bsize); r.font.color.rgb = INK; r.font.name="Calibri"
            if j == 0: r.font.bold = True
    return gt

def callout(s, x, y, w, h, label, text, accent=GREEN, fill=TINT, tcol=INK):
    rrect(s, x, y, w, h, fill, radius=0.08)
    circle(s, x+0.34, y+h/2, 0.16, accent)
    tb(s, x+0.62, y+0.12, w-0.8, h-0.24,
       [[(label+"  ", 11.5, accent, True, False), (text, 11.5, tcol, False, False)]],
       anchor=MSO_ANCHOR.MIDDLE, leading=1.05)

def icon_card(s, x, y, w, h, letter, icol, title, desc):
    rrect(s, x, y, w, h, TINT, line=LINE, lw=1.0, radius=0.07)
    circle(s, x+0.46, y+0.52, 0.56, icol)
    tb(s, x+0.18, y+0.30, 0.56, 0.45, [[(letter, 16, WHITE, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    tb(s, x+0.92, y+0.22, w-1.05, 0.4, [[(title, 14, INK, True, False)]], wrap=False)
    tb(s, x+0.92, y+0.62, w-1.10, h-0.7, [[(desc, 11, MUTE, False, False)]], leading=1.05)

def legend(s, x, y, w, h, items, size=11.5, gap=3, leading=1.04, dark=False):
    """items: list of (tag, color, text)."""
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    tcol = ICE if dark else INK
    for i, (tag, color, text) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.space_before = Pt(0); p.line_spacing = leading
        rn = p.add_run(); rn.text = f"{tag}  "; rn.font.size = Pt(size); rn.font.bold = True
        rn.font.color.rgb = color; rn.font.name = "Calibri"
        rt = p.add_run(); rt.text = text; rt.font.size = Pt(size); rt.font.color.rgb = tcol; rt.font.name = "Calibri"
    return box

# =====================================================================
# S1 — TITLE (dark)
# =====================================================================
s = slide(NAVY)
rrect(s, 0, 0, SW, SH, NAVY)  # ensure full dark
pill(s, 0.75, 0.62, 5.5, 0.46, NAVY2, "AVFinal · UniCEUB · Introdução à Computação · 2026", GOLD, size=11)
t = tb(s, 0.72, 1.35, 11.5, 1.2, [[("OpenBox ", 54, WHITE, True, False), ("v0.2.0", 54, GREEN, True, False)]], wrap=False)
t.text_frame.paragraphs[0].runs[0].font.name = "Cambria"
t.text_frame.paragraphs[0].runs[1].font.name = "Cambria"
tb(s, 0.75, 2.55, 11.4, 0.9, [[("Plataforma de borda com privacidade auditável, hardening e monitoramento integrado", 20, ICE, False, False)]], leading=1.1)
# chips
chips = ["WireGuard","nftables","DNSCrypt","Pi-hole","Unbound","Tor","ADV7Sec"]
cx = 0.75
for c in chips:
    w = 0.34 + len(c)*0.105
    pill(s, cx, 3.45, w, 0.4, NAVY2, c, ICE, size=10.5, bold=False)
    cx += w + 0.18
# hero image
framed_pic(s, "fig1_arch.png", 0.75, 4.25, 11.85, 2.55, dark=True)
tb(s, 0.75, 6.95, 11.85, 0.4, [[("Alcides Olivo Pollazzon Soterio   ·   Brasília–DF · 2026", 11, MUTEDK, False, False)]], wrap=False)

# =====================================================================
# S2 — AGENDA (light)
# =====================================================================
s = slide(WHITE)
y0 = header(s, "≡", "Agenda", "O que vamos percorrer")
agenda = ["Problema e contexto","Estudo de caso: TV box","Linha do tempo","Arquitetura de borda",
          "Fluxo de tráfego","Modelo de ameaça e risco","Plano de hardening","ADV7Sec — detecção",
          "Runbook e evidências","Conclusão e carreira"]
colx = [0.9, 4.55]
for i, item in enumerate(agenda):
    col = i // 5; row = i % 5
    x = colx[col]; y = y0 + 0.15 + row*0.86
    circle(s, x+0.26, y+0.26, 0.52, NAVY)
    tb(s, x, y+0.02, 0.52, 0.5, [[(f"{i+1:02d}", 12.5, GOLD, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    tb(s, x+0.72, y+0.04, 3.0, 0.5, [[(item, 13.5, INK, False, False)]], anchor=MSO_ANCHOR.MIDDLE, leading=1.0)
# right visual panel
framed_pic(s, "thumb_risk.png", 8.55, 2.35, 4.1, 3.0, dark=False)
tb(s, 8.55, 5.45, 4.1, 0.4, [[("Da ameaça de fábrica à borda auditável", 11.5, MUTE, False, True)]], align=PP_ALIGN.CENTER)
tb(s, 0.9, 6.45, 7.4, 0.5, [[("+ Anexo técnico  ", 12.5, GOLD, True, False),
   ("teardown de hardware (R3290_V8.1) e fase de flash/SSH", 12.5, MUTE, False, False)]], leading=1.0)
footer(s, 2)

# =====================================================================
# S3 — PROBLEMA & CONTEXTO (light)
# =====================================================================
s = slide(WHITE)
y0 = header(s, "01", "Problema", "Dois vetores de risco no lar técnico")
tb(s, 0.9, y0+0.05, 11.6, 0.6, [[("TV boxes genéricas saem de fábrica com firmware potencialmente comprometido; "
      "a workstation de desenvolvimento opera sem trilha forense. São problemas independentes.", 13, MUTE, False, False)]], leading=1.1)
icon_card(s, 0.9, 2.85, 5.7, 1.55, "B", BLUE, "Borda (rede doméstica)",
          "DNS em texto claro pelo ISP; VPN sem kill switch — qualquer queda reintroduz o perfil bruto de tráfego.")
icon_card(s, 6.95, 2.85, 5.7, 1.55, "E", BORDO, "Endpoint (workstation)",
          "Sem baseline de integridade e sem alerta comportamental; o primeiro sinal de comprometimento costuma ser tardio.")
# stats
stat(s, 0.95, 4.85, 3.6, "10M+", "dispositivos infectados — BadBox 2.0", ncol=BORDO)
stat(s, 4.95, 4.85, 3.6, "1.3M", "caixas Vo1d — Brasil é o nº 1", ncol=BORDO)
stat(s, 8.95, 4.85, 3.6, "200–400", "R$/mês em assinaturas substituíveis", ncol=GREEN)
callout(s, 0.9, 6.05, 11.7, 0.78, "Tese.",
        "Segurança tratada de forma sistêmica — hardware, SO, rede, monitoramento, integridade e resposta — não como “antivírus”.", accent=GREEN)
footer(s, 3)

# =====================================================================
# S4 — ESTUDO DE CASO (light)
# =====================================================================
s = slide(WHITE)
y0 = header(s, "02", "Estudo de caso", "MXQ Pro 5G 4K — a ameaça de fábrica")
table(s, 0.9, y0+0.15, 6.4, ["Elemento","Observação"], [
    ["SoC","Rockchip RK3229 (placa R29_5G_LP3)"],
    ["Memória","1 GB LPDDR3 — tuning conservador"],
    ["Ethernet","100 Mbps — coerente com o custo"],
    ["Firmware","comprometido de fábrica em lotes"],
    ["Risco","backdoor, OTA, ADB aberto"],
], colw=[1.7, 4.7], rowh=0.46)
# right: threat + mitigation
rrect(s, 7.65, y0+0.15, 5.0, 2.45, TINT, line=LINE, lw=1.0, radius=0.06)
tb(s, 7.95, y0+0.35, 4.5, 0.4, [[("Cadeia de suprimento", 14, BORDO, True, False)]])
bullets(s, 7.95, y0+0.82, 4.45, 1.7, [
    "BadBox 2.0 / Vo1d infectam o firmware antes da 1ª conexão",
    "Mitigação: wipe via maskrom + flash Armbian auditável",
    "Verificação pós-flash antes de qualquer tráfego de rede",
], size=11.5, accent=BORDO, gap=6)
stat(s, 7.95, y0+2.85, 4.5, "antes do 1º boot", "o problema existe na cadeia, não na rede", ncol=GREEN, nsize=22)
callout(s, 0.9, 6.05, 11.7, 0.78, "Valor acadêmico.",
        "O dispositivo conecta computação básica (ARM, bootchain, eMMC) a segurança real (firmware pré-infectado).", accent=BLUE)
footer(s, 4)

# =====================================================================
# S5 — LINHA DO TEMPO (dark)
# =====================================================================
s = slide(NAVY)
rrect(s, 0, 0, SW, SH, NAVY)
header(s, "03", "Evolução", "Linha do tempo do projeto", dark=True)
ny = 3.55
xs = [1.7, 4.35, 7.0, 9.65, 12.0]
ms = [("TVBox OSS v2","base inicial OSS",BORDO,1),
      ("OpenBox v0.1.0","skeleton da pilha\nde privacidade",BLUE,2),
      ("OpenBox v0.2.0","retarget RK3229 ·\nArmbian · Jellyfin",GREEN,3),
      ("ADV7Box","consolidação\nHTML + PDF",BLUE,4),
      ("AVAL01-IC","pass acadêmico\nAVFinal",GOLD,5)]
cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(xs[0]), Inches(ny), Inches(xs[-1]), Inches(ny))
cn.line.color.rgb = RGBColor(0x35,0x50,0x6B); cn.line.width = Pt(2.5); cn.shadow.inherit=False
for (name,desc,col,num),nx in zip(ms,xs):
    hi = (num==3); d = 0.5 if hi else 0.34
    if hi: circle(s, nx, ny, d+0.2, NAVY, line=GREEN, lw=2.0)
    circle(s, nx, ny, d, col, line=WHITE, lw=1.75)
    tb(s, nx-0.3, ny-0.2, 0.6, 0.4, [[(str(num), 13 if hi else 11, WHITE, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    tb(s, nx-1.28, 2.5, 2.56, 0.62, [[(name, 14, WHITE, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
    tb(s, nx-1.25, ny+0.4, 2.5, 0.95, [[(desc, 11, MUTEDK, False, False)]], align=PP_ALIGN.CENTER, leading=1.0)
callout(s, 0.9, 5.7, 11.7, 0.85, "Da prova de conceito à entrega.",
        "A troca de RPi 4 por RK3229 provou que hardware reaproveitado de baixo custo sustenta uma pilha de segurança honesta.",
        accent=GREEN, fill=NAVY2, tcol=ICE)
footer(s, 5, dark=True)

# =====================================================================
# S6 — ARQUITETURA (light)
# =====================================================================
s = slide(WHITE)
y0 = header(s, "04", "Arquitetura", "Defesa em profundidade, dois pilares")
framed_pic(s, "fig1_arch.png", 0.9, y0+0.1, 11.7, 1.9, dark=False)
ay = y0+2.25
rrect(s, 0.9, ay, 5.7, 2.35, TINT, line=LINE, lw=1.0, radius=0.06)
tb(s, 1.2, ay+0.18, 5.1, 0.4, [[("OpenBox — a borda", 15, NAVY, True, False)]])
bullets(s, 1.2, ay+0.66, 5.15, 1.6, [
    "WireGuard como única via de saída permitida",
    "nftables: política + kill switch atômico (fwmark 51820)",
    "DNS local e cifrado: dnscrypt-proxy + Unbound + Pi-hole",
], size=12, accent=GREEN, gap=6)
rrect(s, 6.9, ay, 5.7, 2.35, TINT, line=LINE, lw=1.0, radius=0.06)
tb(s, 7.2, ay+0.18, 5.1, 0.4, [[("ADV7Sec — o host", 15, BORDO, True, False)]])
bullets(s, 7.2, ay+0.66, 5.15, 1.6, [
    "auditd, Falco, Suricata e Zeek agregados localmente",
    "Plano de detecção — observa, nunca bloqueia (preview-first)",
    "Logs em journald, sem SIEM e sem sink externo",
], size=12, accent=BORDO, gap=6)
footer(s, 6)

# =====================================================================
# S7 — FLUXO DE TRÁFEGO (light)
# =====================================================================
s = slide(WHITE)
y0 = header(s, "05", "Fluxo", "Do cliente à Internet, sem vazar")
framed_pic(s, "fig3_flow.png", 0.9, y0+0.1, 11.7, 2.25, dark=False)
steps = ["Cliente solicita acesso","OpenBox intercepta o DNS","dnscrypt-proxy + Unbound resolvem",
         "WireGuard encapsula o tráfego","nftables bloqueia se a VPN cair","ADV7Sec coleta e alerta"]
sy = y0+2.7
for i, st in enumerate(steps):
    col = i % 3; row = i // 3
    x = 0.95 + col*4.05; y = sy + row*0.62
    circle(s, x+0.18, y+0.18, 0.36, NAVY)
    tb(s, x, y+0.0, 0.36, 0.38, [[(str(i+1), 11, GOLD, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    tb(s, x+0.5, y+0.0, 3.45, 0.4, [[(st, 11.5, INK, False, False)]], anchor=MSO_ANCHOR.MIDDLE, leading=1.0)
callout(s, 0.9, 6.18, 11.7, 0.7, "Critério.",
        "Se a rede cai, nada sai. O kill switch é atômico e mantém allowlist apenas para o endpoint da VPN.", accent=GREEN)
footer(s, 7)

# =====================================================================
# S8 — MODELO DE AMEAÇA & RISCO (light)
# =====================================================================
s = slide(WHITE)
y0 = header(s, "06", "Risco", "Modelo de ameaça e matriz de risco")
table(s, 0.9, y0+0.15, 6.3, ["Falha","Categoria","Nota"], [
    ["Queima do gateway / energia","Hardware","Alto"],
    ["Erro de config. do firewall","Software","Alto"],
    ["DDoS / interceptação","Rede","Alto"],
    ["Vazamento de DNS","Rede","Médio/Alto"],
], colw=[3.5, 1.7, 1.1], rowh=0.5)
framed_pic(s, "fig4_riskmatrix.png", 7.5, y0+0.1, 5.15, 2.7, dark=False)
callout(s, 0.9, 5.95, 11.7, 0.85, "3 categorias, mitigação específica.",
        "Hardware, software e rede mantidas distintas para priorizar controles — não respostas genéricas. "
        "Adversário estatal e acesso físico ficam fora do escopo.", accent=BLUE)
footer(s, 8)

# =====================================================================
# S9 — HARDENING (light)
# =====================================================================
s = slide(WHITE)
y0 = header(s, "07", "Mitigação", "Plano de hardening")
gx, gy, gw, gh, ggap = 0.9, y0+0.2, 5.75, 1.85, 0.35
icon_card(s, gx, gy, gw, gh, "H", BORDO, "Hardware",
          "UPS, backup de configuração, imagem de recuperação e equipamento reserva.")
icon_card(s, gx+gw+ggap, gy, gw, gh, "S", BLUE, "Software",
          "Git e rollback, validação de configuração e testes automatizados antes e depois.")
icon_card(s, gx, gy+gh+ggap, gw, gh, "R", GREEN, "Rede",
          "Firewall, VLAN, VPN obrigatória, DNS criptografado e rate limiting.")
icon_card(s, gx+gw+ggap, gy+gh+ggap, gw, gh, "I", GOLD, "Integridade",
          "AIDE e RKHunter, logs e alertas por ntfy sobre arquivos críticos.")
callout(s, 0.9, 6.35, 11.7, 0.7, "Sem composições frágeis.",
        "Não se mistura UFW com nftables, nem se assume que o Tor sustenta streaming — Tor só para metadados.", accent=GREEN)
footer(s, 9)

# =====================================================================
# S10 — ADV7SEC (light)
# =====================================================================
s = slide(WHITE)
y0 = header(s, "08", "Detecção", "ADV7Sec — telemetria host-side")
table(s, 0.9, y0+0.15, 7.3, ["Sensor","O que observa","Saída"], [
    ["auditd","syscalls e arquivos sensíveis","log estruturado"],
    ["Falco","processos e rede","alerta + JSON"],
    ["Suricata","assinaturas de rede","EVE JSON"],
    ["Zeek","sessões e metadados","conn.log"],
    ["AIDE","integridade de arquivos","baseline / diff"],
], colw=[1.6, 3.7, 2.0], rowh=0.44)
rrect(s, 8.55, y0+0.15, 4.1, 2.2, TINT, line=LINE, lw=1.0, radius=0.06)
tb(s, 8.85, y0+0.35, 3.6, 0.4, [[("Preview-first", 14, GREEN, True, False)]])
bullets(s, 8.85, y0+0.8, 3.55, 1.5, [
    "Mostra o que faria antes de agir",
    "Detecção, não bloqueio ativo",
    "Favorece estudo e depuração",
], size=11.5, accent=GREEN, gap=6)
callout(s, 0.9, 5.95, 11.7, 0.78, "Útil quando a ameaça passou da borda.",
        "Visão coerente de comportamento anômalo no endpoint, com normalização e alertas locais.", accent=BORDO)
footer(s, 10)

# =====================================================================
# S11 — RUNBOOK & EVIDÊNCIAS (light)
# =====================================================================
s = slide(WHITE)
y0 = header(s, "09", "Operação", "Runbook, validação e evidências")
tb(s, 0.9, y0+0.05, 5.8, 0.4, [[("Healthchecks", 14, NAVY, True, False)]])
checks = ["systemctl status wg-quick@wg0 …","curl -4 https://ifconfig.io","openbox-dnsleak-check.sh","lynis audit system --quick"]
cb = s.shapes.add_textbox(Inches(0.9), Inches(y0+0.5), Inches(5.8), Inches(2.0))
ctf = cb.text_frame; ctf.word_wrap=True; ctf.margin_left=0; ctf.margin_top=0; ctf.margin_right=0; ctf.margin_bottom=0
for i, cmd in enumerate(checks):
    p = ctf.paragraphs[0] if i==0 else ctf.add_paragraph()
    p.space_after = Pt(8); p.line_spacing=1.0
    rr = p.add_run(); rr.text = "$ "+cmd; rr.font.name="Consolas"; rr.font.size=Pt(11); rr.font.color.rgb=INK
tb(s, 7.1, y0+0.05, 5.5, 0.4, [[("Artefatos versionados", 14, NAVY, True, False)]])
bullets(s, 7.1, y0+0.5, 5.5, 2.2, [
    "README.md — visão geral e diferenças do TVBox OSS",
    "CHANGELOG.md — histórico v0.1.0 → v0.2.0",
    "THREAT_MODEL.md — adversários e fronteiras",
    "RK3229_THREAT_RESEARCH.md — cadeia de suprimento",
    "RUNBOOK.md — operação e troubleshooting",
    "OBSIDIAN_VAULT.md — trilha de execução",
], size=12, accent=BLUE, gap=5)
stat(s, 0.95, y0+2.7, 5.7, "14 págs · 3 formatos · 5 figuras", "documentação reprodutível e auditável", ncol=GREEN, nsize=20)
footer(s, 11)

# =====================================================================
# S12 — CONCLUSÃO & CARREIRA (dark)
# =====================================================================
s = slide(NAVY)
rrect(s, 0, 0, SW, SH, NAVY)
header(s, "10", "Fechamento", "Borda mais segura, auditável e replicável", dark=True)
tb(s, 0.9, 2.1, 7.0, 2.4, [
    [("Direcionamento de carreira", 16, GOLD, True, False)],
    [("Forte viés ", 14, ICE, False, False), ("Blue Team", 14, GREEN, True, False),
     (": hardening, observabilidade, resposta a incidentes e documentação técnica.", 14, ICE, False, False)],
    [("Compatível com SOC, administração de redes, segurança de sistemas e engenharia de plataformas.", 14, ICE, False, False)],
    [("Evolui para automação maior, análise de alertas e inteligência local em semestres futuros.", 14, MUTEDK, False, False)],
], sp=8, leading=1.12)
framed_pic(s, "fig5_closing.png", 8.2, 2.15, 4.4, 1.6, dark=True)
stat(s, 8.4, 4.2, 4.2, "< R$ 1.000", "custo de entrada do laboratório", ncol=GREEN, nsize=30)
callout(s, 0.9, 5.85, 11.7, 0.95, "Em termos acadêmicos e práticos.",
        "Estuda um cenário real, justifica a tecnologia e demonstra defesa — entregando uma borda residencial "
        "compreensível e verificável.", accent=GREEN, fill=NAVY2, tcol=ICE)
tb(s, 0.9, 6.95, 11.7, 0.4, [[("Obrigado.  ", 13, WHITE, True, False),
    ("Alcides Olivo Pollazzon Soterio · UniCEUB · Introdução à Computação · 2026", 11, MUTEDK, False, False)]], wrap=False)

# =====================================================================
# A1 — ANEXO: estruturas (face superior) — light
# =====================================================================
s = slide(WHITE)
header(s, "A1", "Anexo · Hardware", "Estruturas da placa — face superior")
framed_pic(s, "hw/plate_top.png", 0.7, 2.0, 3.9, 4.95, dark=False)
legend(s, 5.0, 2.05, 7.8, 4.9, [
    ("1", BLUE,  "eMMC Samsung KLM8G1GETF — 8 GB (armazenamento)"),
    ("2", RGBColor(0xB5,0x65,0x1D), "R-série — resistores no barramento eMMC (tap)"),
    ("3", RGBColor(0xB5,0x65,0x1D), "Trilha CLK (breakout) — short p/ maskrom"),
    ("4", GREEN, "QFN-44 SVS6Z56F / TAC2121 — PHY/interface (provável)"),
    ("5", GREEN, "PPT PM44-11BP — driver/MOSFET de potência (provável)"),
    ("6", GOLD,  "Capacitor eletrolítico CK 100 µF — filtro de entrada"),
    ("7", GOLD,  "LDO AMS1117 (1117C) — regulador 3,3 V"),
    ("8", GOLD,  "Buck 3R3 (3,3 µH) ×2 — trilhos de núcleo/DDR"),
    ("9", NAVY,  "SoC Rockchip RK3229 — sob o dissipador (família rk322x)"),
    ("10", BORDO,"Header serial R / T / G — console UART de depuração"),
], size=13, gap=7)
footer(s, 13)

# =====================================================================
# A2 — ANEXO: memória, face inferior e portas — light
# =====================================================================
s = slide(WHITE)
header(s, "A2", "Anexo · Hardware", "Memória, face inferior e portas")
framed_pic(s, "hw/plate_bottom.png", 0.7, 2.0, 5.05, 2.5, dark=False)
framed_pic(s, "hw/plate_ports.png", 0.7, 4.72, 5.05, 2.3, dark=False)
legend(s, 6.1, 2.05, 6.7, 5.0, [
    ("11–12", BLUE, "DRAM Hynix H5TQ2G43BFR (DDR3) — 4 chips ≈ 1 GB"),
    ("13", GOLD,  "Rede de resistores — terminação/strap do DDR"),
    ("14", GOLD,  "LDO 1117 (face inferior) — regulador auxiliar"),
    ("15", BLUE,  "Footprint BGA vazio — posição NAND/eMMC alternativa"),
    ("16", RGBColor(0xB5,0x65,0x1D), "Cristal/oscilador — referência de clock"),
    ("17", GREEN, "USB-A — porta host USB 2.0"),
    ("18–19", GREEN, "AV 3,5 mm — saídas CVBS/áudio"),
    ("20", GREEN, "HDMI — saída de vídeo digital"),
    ("21", GREEN, "RJ45 — Ethernet 10/100"),
    ("22", GOLD,  "DC jack — entrada de alimentação 5 V"),
], size=12.5, gap=6)
callout(s, 6.1, 6.25, 6.7, 0.85, "Privacidade & flash.",
        "MAC redigido na face inferior; o micro-USB OTG é a porta usada para o flash via maskrom.", accent=BORDO)
footer(s, 14)

# =====================================================================
# A3 — ANEXO: fase de flash e SSH — dark
# =====================================================================
s = slide(NAVY)
rrect(s, 0, 0, SW, SH, NAVY)
header(s, "A3", "Anexo · Procedimento", "Fase de flash e SSH (rk322x)", dark=True)
framed_pic(s, "hw/term_flash.png", 0.55, 2.05, 6.25, 3.2, dark=True)
framed_pic(s, "hw/term_ssh.png", 6.95, 2.05, 6.0, 3.6, dark=True)
callout(s, 0.9, 6.1, 11.95, 0.92, "Reconstrução do procedimento documentado.",
        "Valores reais do projeto — placa R3290_V8.1 · Armbian rk322x-box · host 192.168.0.103 · Phase 0 PASS (4/4). "
        "Sequência: maskrom → rkdeveloptool → boot → SSH por chave.", accent=GREEN, fill=NAVY2, tcol=ICE)
footer(s, 15, dark=True)

out = HERE / "OpenBox_Apresentacao.pptx"
prs.save(str(out))
print("DECK saved:", out, out.stat().st_size, "bytes ·", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
