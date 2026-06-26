# -*- coding: utf-8 -*-
"""OpenBox project timeline slide (recap) with snapshot figures."""
import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn

HERE = pathlib.Path(__file__).parent
FIGS = HERE / "figs"

# palette — midnight navy dominant, blue/green support, gold accent
NAVY   = RGBColor(0x0C, 0x22, 0x33)   # background
SLATE  = RGBColor(0x35, 0x50, 0x6B)   # timeline line
BLUE   = RGBColor(0x2E, 0x6F, 0xA8)
GREEN  = RGBColor(0x1F, 0x9E, 0x75)
BORDO  = RGBColor(0x8A, 0x1F, 0x3D)
GOLD   = RGBColor(0xF2, 0xB7, 0x05)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ICE    = RGBColor(0xCA, 0xDC, 0xFC)
MUTE   = RGBColor(0x8F, 0xA3, 0xB8)
CARDBD = RGBColor(0x24, 0x40, 0x5C)
PILL   = RGBColor(0x14, 0x31, 0x4C)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# background
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = NAVY

def textbox(x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    """lines: list of (text, size, color, bold, italic, spacing_after_pt)"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, (text, size, color, bold, italic, sp) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if sp is not None:
            p.space_after = Pt(sp)
        p.space_before = Pt(0)
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = "Calibri"
    return tb

def circle(cx, cy, d, fill, line=None, line_w=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx-d/2), Inches(cy-d/2), Inches(d), Inches(d))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s

def rounded(x, y, w, h, fill, line=None, line_w=1.0, radius=0.08):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    try:
        s.adjustments[0] = radius
    except Exception:
        pass
    return s

# ---------- title ----------
t = textbox(0.6, 0.42, 9.6, 0.9, [("OpenBox — Linha do Tempo do Projeto", 34, WHITE, True, False, 0)])
t.text_frame.paragraphs[0].runs[0].font.name = "Cambria"
textbox(0.62, 1.24, 9.7, 0.7,
        [("Da base TVBox OSS à entrega acadêmica AVFinal — privacidade de borda auditável em Rockchip RK3229",
          15, ICE, False, False, 0)])

# tag pill top-right
pill = rounded(10.45, 0.52, 2.45, 0.5, PILL, radius=0.5)
tp = pill.text_frame; tp.word_wrap = False
tp.margin_left = 0; tp.margin_right = 0; tp.margin_top = 0; tp.margin_bottom = 0
p = tp.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "AVFinal · UniCEUB · 2026"
r.font.size = Pt(11.5); r.font.bold = True; r.font.color.rgb = GOLD; r.font.name = "Calibri"

# ---------- timeline ----------
ny = 2.92
xs = [1.55, 4.13, 6.70, 9.28, 11.85]
milestones = [
    ("TVBox OSS v2", "base inicial OSS", BORDO, 1),
    ("OpenBox v0.1.0", "skeleton: WireGuard,\nnftables, Pi-hole", BLUE, 2),
    ("OpenBox v0.2.0", "retarget RK3229 ·\nArmbian · Jellyfin", GREEN, 3),
    ("ADV7Box", "consolidação\nHTML + PDF", BLUE, 4),
    ("AVAL01-IC", "pass acadêmico\nAVFinal", GOLD, 5),
]
# baseline connector
conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(xs[0]), Inches(ny), Inches(xs[-1]), Inches(ny))
conn.line.color.rgb = SLATE; conn.line.width = Pt(2.5)
conn.shadow.inherit = False

for (name, desc, col, num), nx in zip(milestones, xs):
    highlight = (num == 3)
    d = 0.52 if highlight else 0.34
    if highlight:
        circle(nx, ny, d+0.20, NAVY, line=GREEN, line_w=2.0)   # halo ring
    circle(nx, ny, d, col, line=WHITE, line_w=1.75)
    # number inside node
    nb = textbox(nx-0.3, ny-0.20, 0.6, 0.4, [(str(num), 13 if highlight else 11, WHITE, True, False, 0)],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    # version label above
    textbox(nx-1.25, 1.95, 2.5, 0.62, [(name, 13.5, WHITE, True, False, 0)],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
    # one-liner below
    textbox(nx-1.22, 3.30, 2.44, 0.95, [(desc, 10.5, MUTE, False, False, 0)],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

# ---------- snapshots band ----------
textbox(0.6, 4.30, 6.0, 0.4, [("Snapshots do projeto", 14, GOLD, True, False, 0)])

cards = [
    ("thumb_arch.png", "Arquitetura — defesa em profundidade"),
    ("thumb_flow.png", "Fluxo de tráfego — 6 etapas"),
    ("thumb_risk.png", "Matriz de risco — HW · SW · Rede"),
]
card_w = 3.8
gap = 0.36
x0 = 0.6
img_w = 3.3
img_h = img_w / (1000/600)   # thumbnails are 1000x600
img_y = 4.74
for i, (img, cap) in enumerate(cards):
    cx = x0 + i * (card_w + gap)
    img_x = cx + (card_w - img_w) / 2
    # frame
    rounded(img_x - 0.06, img_y - 0.06, img_w + 0.12, img_h + 0.12, WHITE, line=CARDBD, line_w=1.0, radius=0.04)
    slide.shapes.add_picture(str(FIGS / img), Inches(img_x), Inches(img_y), Inches(img_w), Inches(img_h))
    # caption
    textbox(cx, img_y + img_h + 0.12, card_w, 0.5,
            [(cap, 11, ICE, False, False, 0)], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

out = HERE / "OpenBox_Timeline.pptx"
prs.save(str(out))
print("PPTX saved:", out, out.stat().st_size, "bytes")
print("img_h =", round(img_h,3), "caption_y =", round(img_y+img_h+0.12,3))
